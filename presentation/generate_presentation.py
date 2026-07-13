"""Generate the 10-slide Rossmann executive PowerPoint."""
from pathlib import Path
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT=Path(__file__).resolve().parents[1]
NAVY=RGBColor(17,31,58); DARK=RGBColor(35,41,50); MUTED=RGBColor(91,103,120)

def ensure_png(name):
    path=ROOT/"visuals"/name
    if not path.exists():
        svg=path.with_suffix('.svg')
        if svg.exists():
            from cairosvg import svg2png
            svg2png(url=str(svg),write_to=str(path))
    return path

def title(slide,heading,subtitle=''):
    box=slide.shapes.add_textbox(Inches(.7),Inches(.4),Inches(12),Inches(.75)); p=box.text_frame.paragraphs[0]; p.text=heading; p.font.size=Pt(27); p.font.bold=True; p.font.color.rgb=NAVY
    if subtitle:
        box=slide.shapes.add_textbox(Inches(.72),Inches(1.08),Inches(11.8),Inches(.4)); p=box.text_frame.paragraphs[0]; p.text=subtitle; p.font.size=Pt(11); p.font.color.rgb=MUTED

def bullets(slide,items,x=.8,y=1.55,w=5.7,h=4.9,size=18):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear()
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.size=Pt(size); p.font.color.rgb=DARK; p.space_after=Pt(10)

def image(slide,name,x=6.7,y=1.55,w=5.8,h=4.8):
    path=ensure_png(name)
    if path.exists(): slide.shapes.add_picture(str(path),Inches(x),Inches(y),width=Inches(w),height=Inches(h))

def main():
    metrics_path=ROOT/"reports/benchmark_metrics.json"; metrics=json.loads(metrics_path.read_text()) if metrics_path.exists() else {"recommended_model":"XGBoost"}
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    slide=prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=NAVY
    box=slide.shapes.add_textbox(Inches(.8),Inches(1.15),Inches(11.7),Inches(1.2)); p=box.text_frame.paragraphs[0]; p.text='Rossmann Store Sales'; p.font.size=Pt(40); p.font.bold=True; p.font.color.rgb=RGBColor(255,255,255)
    box=slide.shapes.add_textbox(Inches(.82),Inches(2.55),Inches(11),Inches(1)); p=box.text_frame.paragraphs[0]; p.text='Six-week retail forecasting, leakage-safe validation, and business planning'; p.font.size=Pt(22); p.font.color.rgb=RGBColor(210,222,241)
    specs=[
        ('2. Business Problem','Forecast demand to coordinate inventory, staffing, promotions, and cash flow',['Under-forecasting creates stockouts and lost sales.','Over-forecasting ties up inventory and labor.','The operating horizon is six weeks across 1,115 stores.'],'monthly_sales.png'),
        ('3. Dataset','Kaggle Rossmann Store Sales competition',['1,017,209 training store-days','41,088 public test rows','Sales is the target','Customers is excluded because it is unavailable at inference'],'sales_distribution.png'),
        ('4. Exploratory Analysis','Store, calendar, promotion, and closure effects drive sales',['Sales are right-skewed and seasonal.','Promotions lift demand.','Closed stores have zero sales.','Store metadata contains structural missingness.'],'promo_lift.png'),
        ('5. Modeling','Transparent baseline versus nonlinear boosting',['Ridge Regression baseline','XGBoost advanced model','log1p target','42-day validation and 42-day test windows'],'correlation_heatmap.png'),
        ('6. Results','RMSPE is the primary competition metric',[f"Recommended demo model: {metrics.get('recommended_model','XGBoost')}",'RMSE, MAE, MAPE, and R² provide supporting context.','Final metrics must be regenerated from authenticated Kaggle files.'],'model_comparison.png'),
        ('7. Insights','Forecasting quality is heterogeneous',['Promotion-day errors require separate monitoring.','High-volume stores dominate absolute error.','Holiday periods and closures need explicit rules.'],'sales_by_day.png'),
        ('8. Recommendations','Use forecasts as governed planning inputs',['Combine with safety stock and lead times.','Review high-error stores weekly.','Track bias and RMSPE by segment.','Use human override for structural changes.'],None),
        ('9. Future Work','Advance from point forecasts to decision intelligence',['Prediction intervals and quantile loss','Hierarchical reconciliation','Causal promotion uplift','Weather, events, inventory, and drift monitoring'],None),
    ]
    for heading,sub,items,chart in specs:
        slide=prs.slides.add_slide(blank); title(slide,heading,sub)
        if chart: bullets(slide,items); image(slide,chart)
        else: bullets(slide,items,1.0,1.7,11.2,4.7,21)
    slide=prs.slides.add_slide(blank); slide.background.fill.solid(); slide.background.fill.fore_color.rgb=NAVY
    box=slide.shapes.add_textbox(Inches(.8),Inches(1.2),Inches(11.7),Inches(.9)); p=box.text_frame.paragraphs[0]; p.text='10. Conclusion'; p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=RGBColor(255,255,255)
    box=slide.shapes.add_textbox(Inches(1),Inches(2.45),Inches(11.2),Inches(2.4)); p=box.text_frame.paragraphs[0]; p.text='Reliable retail forecasting combines leakage-safe modeling, chronological evaluation, operational business rules, uncertainty management, and continuous monitoring.'; p.font.size=Pt(27); p.font.color.rgb=RGBColor(220,230,245); p.alignment=PP_ALIGN.CENTER
    out=Path(__file__).with_name('Rossmann_Store_Sales_Executive_Presentation.pptx'); prs.save(out); print(f'Saved {out}')

if __name__=='__main__': main()
